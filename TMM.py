# -*- coding: utf-8 -*-
"""
Created on Thu May 30 14:10:24 2024

@author: mudit
"""

import os
from pptx import Presentation
from comtypes.client import CreateObject
import comtypes

# Define the folder containing the PowerPoint files
folder_path = "C:\\Script4TraininigModules\\Modules"
old_footer_text = "Oct"
new_footer_text = "Dec"
text_to_replace_in_filename = "BB3"
replacement_text_in_filename = "BB5"

# Function to replace text in filenames
def rename_files():
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            new_name = filename.replace(text_to_replace_in_filename, replacement_text_in_filename)
            os.rename(os.path.join(folder_path, filename), os.path.join(folder_path, new_name))
            print(new_name)

# Function to convert PowerPoint files to PDFs
def convert_to_pdf():
    powerpoint = CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            ppt_path = os.path.join(folder_path, filename)
            pdf_path = ppt_path.replace(".pptx", ".pdf")
            deck = powerpoint.Presentations.Open(ppt_path)
            deck.SaveAs(pdf_path, 32)  # 32 is the format type for PDF
            deck.Close()

    powerpoint.Quit()

# Function to find and replace text in the master slides
def replace_text_in_master():
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            ppt_path = os.path.join(folder_path, filename)
            presentation = Presentation(ppt_path)
            for slide_master in presentation.slide_masters:
                for shape in slide_master.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_footer_text in run.text:
                                    run.text = run.text.replace(old_footer_text, new_footer_text)
                for layout in slide_master.slide_layouts:
                    for shape in layout.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if old_footer_text in run.text:
                                        run.text = run.text.replace(old_footer_text, new_footer_text)
            presentation.save(ppt_path)

# Run the functions
#rename_files()
#replace_text_in_master()
convert_to_pdf()

print("Task completed successfully.")
