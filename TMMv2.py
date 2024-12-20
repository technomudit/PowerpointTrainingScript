# -*- coding: utf-8 -*-
"""
Created on Thu May 30 14:10:24 2024

@author: mudit
"""

import os
import sys
from pptx import Presentation
from comtypes.client import CreateObject
import comtypes
import tkinter as tk
from tkinter import filedialog
from tkinter import scrolledtext
from tkinter import PhotoImage

# Define the folder containing the PowerPoint files
folder_path = "C:\\Script4TraininigModules\\Modules"
old_footer_text = "Oct"
new_footer_text = "Dec"
text_to_replace_in_filename = "BB3"
replacement_text_in_filename = "BB5"

# Function to replace text in filenames
def rename_files():
    updated_filenames = []
    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            new_name = filename.replace(text_to_replace_in_filename, replacement_text_in_filename)
            os.rename(os.path.join(folder_path, filename), os.path.join(folder_path, new_name))
            updated_filenames.append(new_name)
            console_output.insert(tk.END, f"Renamed: {new_name}\n")
    return updated_filenames

# Function to convert PowerPoint files to PDFs
def convert_to_pdf(filenames):
    powerpoint = CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    for filename in filenames:
        ppt_path = os.path.join(folder_path, filename)
        pdf_path = ppt_path.replace(".pptx", ".pdf")
        deck = powerpoint.Presentations.Open(ppt_path)
        deck.SaveAs(pdf_path, 32)  # 32 is the format type for PDF
        deck.Close()
        console_output.insert(tk.END, f"Converted to PDF: {pdf_path}\n")

    powerpoint.Quit()

# Function to find and replace text in the master slides
def replace_text_in_master(filenames):
    for filename in filenames:
        ppt_path = os.path.join(folder_path, filename)
        presentation = Presentation(ppt_path)
        for slide_master in presentation.slide_masters:
            for shape in slide_master.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_footer_text in run.text:
                                run.text = run.text.replace(old_footer_text, new_footer_text)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if old_footer_text in run.text:
                                    run.text = run.text.replace(old_footer_text, new_footer_text)
        presentation.save(ppt_path)
        console_output.insert(tk.END, f"Replaced text in: {ppt_path}\n")

# Function to select folder
def select_folder():
    global folder_path
    folder_path = filedialog.askdirectory().replace("/", "\\")
    folder_path_entry.delete(0, tk.END)
    folder_path_entry.insert(0, folder_path)

# Function to run all tasks
def run_tasks():
    global old_footer_text, new_footer_text, text_to_replace_in_filename, replacement_text_in_filename, folder_path
    folder_path = folder_path_entry.get().replace("/", "\\")
    old_footer_text = old_footer_text_entry.get()
    new_footer_text = new_footer_text_entry.get()
    text_to_replace_in_filename = text_to_replace_in_filename_entry.get()
    replacement_text_in_filename = replacement_text_in_filename_entry.get()
    updated_filenames = rename_files()
    replace_text_in_master(updated_filenames)
    if convert_to_pdf_var.get():
        convert_to_pdf(updated_filenames)
    console_output.insert(tk.END, "Task completed successfully.\n")

# GUI setup
root = tk.Tk()
root.title("BlueKei Trainings PowerPoint Automator")

# Add company logo
try:
    if hasattr(sys, '_MEIPASS'):
        logo_path = os.path.join(sys._MEIPASS, 'BK_Final.png')
    else:
        logo_path = 'C:\\Script4TraininigModules\\BK_Final.png'
    logo = PhotoImage(file=logo_path)
    tk.Label(root, image=logo).grid(row=0, column=0, columnspan=3, padx=10, pady=10)
except tk.TclError:
    tk.Label(root, text="BlueKei Solutions Pvt Ltd", fg="red").grid(row=0, column=0, columnspan=3, padx=10, pady=10)

# Instructional text
instructions = """Instructions:
1. Select the folder containing the PowerPoint files.
2. Enter the old and new footer text.
3. Enter the text to replace in filenames and the replacement text.
4. Check the 'Convert to PDF' option if you want to convert the files to PDF.
5. Click 'Run' to execute the tasks.

Author: Mudit Mittal
Version: V1.0
Date: 19 Dec 2024
BlueKei Solutions Pvt. Ltd. (c) 2024
"""
tk.Label(root, text=instructions, justify=tk.LEFT).grid(row=1, column=0, columnspan=3, padx=10, pady=10, sticky=tk.W)

tk.Label(root, text="Folder Path:", anchor="w").grid(row=2, column=0, padx=10, pady=5, sticky=tk.W)
folder_path_entry = tk.Entry(root, width=50)
folder_path_entry.grid(row=2, column=1, padx=10, pady=5, sticky=tk.W)
folder_path_entry.insert(0, folder_path)  # Set default folder path
browse_button = tk.Button(root, text="Browse", command=select_folder, fg="blue")
browse_button.grid(row=2, column=2, padx=10, pady=5, sticky=tk.W)

tk.Label(root, text="Old Footer Text:", anchor="w").grid(row=3, column=0, padx=10, pady=5, sticky=tk.W)
old_footer_text_entry = tk.Entry(root, width=50)
old_footer_text_entry.grid(row=3, column=1, padx=10, pady=5, sticky=tk.W)

tk.Label(root, text="New Footer Text:", anchor="w").grid(row=4, column=0, padx=10, pady=5, sticky=tk.W)
new_footer_text_entry = tk.Entry(root, width=50)
new_footer_text_entry.grid(row=4, column=1, padx=10, pady=5, sticky=tk.W)

tk.Label(root, text="Text to Replace in Filename:", anchor="w").grid(row=5, column=0, padx=10, pady=5, sticky=tk.W)
text_to_replace_in_filename_entry = tk.Entry(root, width=50)
text_to_replace_in_filename_entry.grid(row=5, column=1, padx=10, pady=5, sticky=tk.W)

tk.Label(root, text="Replacement Text in Filename:", anchor="w").grid(row=6, column=0, padx=10, pady=5, sticky=tk.W)
replacement_text_in_filename_entry = tk.Entry(root, width=50)
replacement_text_in_filename_entry.grid(row=6, column=1, padx=10, pady=5, sticky=tk.W)

convert_to_pdf_var = tk.BooleanVar()
convert_to_pdf_var.set(True)
tk.Checkbutton(root, text="Convert to PDF", variable=convert_to_pdf_var).grid(row=7, column=1, padx=10, pady=5, sticky=tk.W)

run_button = tk.Button(root, text="Run", command=run_tasks, fg="blue")
run_button.grid(row=7, column=2, padx=10, pady=5, sticky=tk.W)

console_output = scrolledtext.ScrolledText(root, width=70, height=10)
console_output.grid(row=8, column=0, columnspan=3, padx=10, pady=10, sticky=tk.W)

root.mainloop()

print("Task completed successfully.")
